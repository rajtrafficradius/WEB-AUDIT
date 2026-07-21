"""Contract tests for the pure render engines (xlsx, pptx, deck pdf, markdown)."""

from __future__ import annotations

import re
from pathlib import Path
from typing import Any

import pytest
from openpyxl import load_workbook
from pptx import Presentation
from pptx.util import Inches
from pypdf import PdfReader

from exporters.markdown_summary import render_markdown
from exporters.pptx_deck import render_deck, render_deck_pdf, slide_layouts_for
from exporters.xlsx_workbooks import EMPTY_MESSAGE, NO_CHANGE, render_workbooks

pytestmark = pytest.mark.render

CORE_WORKBOOKS = {
    "Technical_Audit_Report.xlsx",
    "Content_Audit_Workbook.xlsx",
    "Backlink_Audit_Report.xlsx",
    "Competitor_Landscape_Analysis.xlsx",
    "GEO_AEO_Readiness_Scorecard.xlsx",
    "CRO_UX_Findings.xlsx",
    "Tracking_Audit_Report.xlsx",
    "Baseline_Performance_Analysis.xlsx",
    "Master_Keyword_Universe.xlsx",
    "Content_Gap_Analysis.xlsx",
    "Content_Strategy.xlsx",
    "URL_Architecture_Map.xlsx",
    "Cannibalization_Resolution_Plan.xlsx",
    "16_Week_Action_Plan.xlsx",
    "Title_Tag_Optimizations.xlsx",
    "Meta_Description_Optimizations.xlsx",
    "H1_Tags.xlsx",
    "Internal_Link_Map.xlsx",
    "Redirect_Map.xlsx",
    "Canonical_Fixes.xlsx",
    "QC_Report.xlsx",
}
ECOMMERCE_WORKBOOK = "Ecommerce_Audit_Report.xlsx"
LOCAL_WORKBOOK = "GBP_Local_Audit.xlsx"
LINK_WORKBOOKS = {"Referring_Domains.xlsx", "Link_Gap_Opportunities.xlsx"}

# V18 benchmark minimums (sheets, data rows) we must equal or beat.
V18_MINIMUMS: dict[str, tuple[int, int]] = {
    "Technical_Audit_Report.xlsx": (9, 644),
    "Content_Audit_Workbook.xlsx": (5, 953),
    "Backlink_Audit_Report.xlsx": (6, 1172),
    "Master_Keyword_Universe.xlsx": (5, 233),
    "Content_Gap_Analysis.xlsx": (3, 169),
    "GEO_AEO_Readiness_Scorecard.xlsx": (3, 126),
    "Competitor_Landscape_Analysis.xlsx": (3, 37),
    # Deliberately below V18's raw row count: identical same-signal rows are
    # now aggregated into one systemic finding each (quality over volume).
    "CRO_UX_Findings.xlsx": (2, 8),
    "Content_Strategy.xlsx": (3, 52),
    "Ecommerce_Audit_Report.xlsx": (4, 276),
    "Tracking_Audit_Report.xlsx": (2, 15),
    "Baseline_Performance_Analysis.xlsx": (2, 12),
}

_REF = re.compile(r"^[A-Z]+\d+:([A-Z]+)(\d+)$")


def _data_row_count(worksheet: Any) -> int:
    """Number of body rows written under the header (row 4)."""
    match = _REF.match(str(worksheet.auto_filter.ref))
    assert match is not None, f"{worksheet.title} has no autofilter range"
    return int(match.group(2)) - 4


def _sheet_texts(worksheet: Any) -> list[str]:
    texts: list[str] = []
    for row in worksheet.iter_rows(values_only=True):
        texts.extend(str(value) for value in row if value is not None)
    return texts


def _deck_texts(presentation: Any) -> list[str]:
    texts: list[str] = []
    for slide in presentation.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                texts.append(shape.text_frame.text)
    return texts


# --------------------------------------------------------------------------- fixtures


def _page(number: int, **overrides: Any) -> dict[str, Any]:
    url = f"https://acmewidgets.com.au/{'about' if number == 2 else 'shop' if number == 3 else ''}"
    page: dict[str, Any] = {
        "id": f"URL-{number:04d}",
        "evidence_id": f"EV-{number:04d}",
        "original_url": url,
        "original_urls": [url],
        "duplicate_observations": 1,
        "normalized_url": url,
        "status_code": 200,
        "title": f"Acme Widgets page {number}",
        "meta_description": f"Meta description for page {number} of the Acme Widgets site.",
        "h1": f"Heading {number}",
        "canonical_url": url,
        "indexability": "No noindex observed",
        "word_count": 800 + number,
        "internal_links": 20 + number,
        "external_links": 3,
        "redirect_chain": [url],
        "content_type": "text/html; charset=utf-8",
        "body_sha256": f"hash-{number:02d}",
        "page_type": "Homepage" if number == 1 else "Content",
        "links": [url],
        "captured_at": "2026-07-16T02:00:00+00:00",
    }
    page.update(overrides)
    return page


def _market_block() -> dict[str, Any]:
    return {
        "status": "available",
        "provider": "semrush",
        "database": "au",
        "unavailable_reason": None,
        "fetched_at": "2026-07-16T02:10:00+00:00",
        "units_spent": 640,
        "domain": {
            "organic_keywords": 1840,
            "organic_traffic": 9120,
            "organic_cost": 14310.5,
            "adwords_keywords": 62,
            "rank": 412345,
            "authority_score": 34,
            "backlinks_total": 8210,
            "referring_domains": 412,
            "referring_ips": 388,
            "follow_links": 7100,
            "nofollow_links": 1110,
        },
    }


def _keyword(index: int, *, mapped: bool) -> dict[str, Any]:
    stage = ("TOFU", "MOFU", "BOFU")[index % 3]
    return {
        "id": f"KW-{index:04d}",
        "phrase": f"widget phrase {index}",
        "position": (index % 40) + 1 if mapped else None,
        "previous_position": (index % 45) + 1 if mapped else None,
        "search_volume": 40 + index * 3,
        "cpc": round(0.8 + (index % 17) * 0.11, 2),
        "competition": round((index % 90) / 100, 2),
        "results_count": 100000 + index,
        "traffic_share": round((index % 30) / 1000, 4),
        "traffic_cost_share": round((index % 25) / 1000, 4),
        "trend": "0.1,0.2,0.3",
        "landing_url": (
            f"https://acmewidgets.com.au/products/widget-{index % 150}" if mapped else None
        ),
        "intent": ("informational", "commercial", "transactional")[index % 3],
        "funnel_stage": stage,
        "cluster": f"Cluster {index % 24:02d}",
        "page_type": "Product" if mapped else "Editorial",
        "opportunity": "Measured position outside the top 10" if mapped else "No page targets it",
        "evidence_ids": [f"EV-KW-{index:04d}"],
        "source": "semrush",
        "unavailable_reason": None,
    }


def _sample_data() -> dict[str, Any]:
    """Small run-data sample carrying every key of the shared contract."""
    pages = [
        _page(
            1,
            schema_types=["Organization"],
            images_total=10,
            images_missing_alt=2,
            response_ms=350,
            body_bytes=48000,
            analytics_tags=["GA4"],
        ),
        _page(
            2,
            word_count=120,
            canonical_url="https://acmewidgets.com.au/about-us",
            response_ms=90,
            body_bytes=22000,
            analytics_tags=[],
        ),
        _page(
            3,
            status_code=404,
            title=None,
            meta_description=None,
            h1=None,
            canonical_url=None,
            word_count=None,
            redirect_chain=[
                "https://acmewidgets.com.au/old-shop",
                "https://acmewidgets.com.au/shop",
            ],
        ),
    ]
    findings = [
        {
            "id": "F-001", "priority": "P1", "priority_score": 92.0,
            "category": "technical", "rule_id": "technical.http_status",
            "rule_version": "1.0.0", "severity": "Critical",
            "title": "HTTP 404 response on a linked page",
            "description": "A crawled URL returned 404 while still receiving internal links.",
            "impact": "Wasted crawl budget and broken user journeys.",
            "confidence": 1.0, "reach": "1 page", "affected_count": 1,
            "affected_urls": ["https://acmewidgets.com.au/shop"], "effort": "S",
            "implementation_risk": "low", "approval_class": "agency_admin",
            "as_of_date": "2026-07-16", "evidence_ids": ["EV-0003"],
        },
        {
            "id": "F-002", "priority": "P2", "priority_score": 71.0,
            "category": "on_page", "rule_id": "on_page.title_length",
            "rule_version": "1.0.0", "severity": "High",
            "title": "Title tag outside recommended length",
            "description": "One page title exceeds the recommended pixel budget.",
            "impact": "Truncated snippets reduce click-through.",
            "confidence": 0.9, "reach": "1 page", "affected_count": 1,
            "affected_urls": ["https://acmewidgets.com.au/"], "effort": "S",
            "implementation_risk": "low", "approval_class": "editorial",
            "as_of_date": "2026-07-16", "evidence_ids": ["EV-0001"],
        },
        {
            "id": "F-003", "priority": "P3", "priority_score": 45.0,
            "category": "performance", "rule_id": "performance.response_ms",
            "rule_version": "1.0.0", "severity": "Medium",
            "title": "Slow HTML response observed",
            "description": "The homepage HTML response exceeded 300 ms during the crawl.",
            "impact": "Slower first byte delays every downstream metric.",
            "confidence": 0.8, "reach": "1 page", "affected_count": 1,
            "affected_urls": ["https://acmewidgets.com.au/"], "effort": "M",
            "implementation_risk": "medium", "approval_class": "agency_admin",
            "as_of_date": "2026-07-16", "evidence_ids": ["EV-0001"],
        },
        {
            "id": "F-004", "priority": "P4", "priority_score": 20.0,
            "category": "keyword_architecture", "rule_id": "keywords.overlap",
            "rule_version": "1.0.0", "severity": "Low",
            "title": "Two pages target overlapping topics",
            "description": "The about and shop pages share heading topics.",
            "impact": "Potential internal competition for one intent.",
            "confidence": 0.7, "reach": "2 pages", "affected_count": 2,
            "affected_urls": [
                "https://acmewidgets.com.au/about",
                "https://acmewidgets.com.au/shop",
            ],
            "effort": "M", "implementation_risk": "low",
            "approval_class": "editorial", "as_of_date": "2026-07-16",
            "evidence_ids": ["EV-0002"],
        },
    ]
    actions = [
        {
            "id": "A-001", "phase": "Foundation", "week": 1, "week_end": 1,
            "priority": "P1", "action": "Approve the evidence boundary",
            "owner": "Agency admin", "dependencies": [], "effort": "3h",
            "kpi": "Scope decision recorded", "approval_class": "agency_admin",
            "status": "Ready", "evidence_ids": ["EV-0003"], "confidence": 1.0,
            "implementation_risk": "low", "notes": "Advisory only.",
        },
        {
            "id": "A-002", "phase": "Foundation", "week": 2, "week_end": 2,
            "priority": "P1", "action": "Fix the 404 shop route",
            "owner": "Developer", "dependencies": ["A-001"], "effort": "4h",
            "kpi": "Route returns 200", "approval_class": "agency_admin",
            "status": "Ready", "evidence_ids": ["EV-0003"], "confidence": 1.0,
            "implementation_risk": "low", "notes": None,
        },
        {
            "id": "A-003", "phase": "Technical Hardening", "week": 3, "week_end": 6,
            "priority": "P2", "action": "Resolve canonical mismatch on the about page",
            "owner": "Developer", "dependencies": [], "effort": "6h",
            "kpi": "Canonical matches the normalized URL",
            "approval_class": "agency_admin", "status": "Ready",
            "evidence_ids": ["EV-0002"], "confidence": 0.9,
            "implementation_risk": "medium", "notes": None,
        },
        {
            "id": "A-004", "phase": "Content", "week": 7, "week_end": 10,
            "priority": "P2", "action": "Rewrite the oversized homepage title",
            "owner": "Editor", "dependencies": [], "effort": "2h",
            "kpi": "Title within length budget", "approval_class": "editorial",
            "status": "Ready", "evidence_ids": ["EV-0001"], "confidence": 0.9,
            "implementation_risk": "low", "notes": None,
        },
        {
            "id": "A-005", "phase": "Content", "week": 9, "week_end": 12,
            "priority": "P3", "action": "Expand the thin about page",
            "owner": "Editor", "dependencies": [], "effort": "8h",
            "kpi": "Word count above threshold", "approval_class": "editorial",
            "status": "Ready", "evidence_ids": [], "confidence": 0.8,
            "implementation_risk": "low", "notes": None,
        },
        {
            "id": "A-006", "phase": "Measurement", "week": 13, "week_end": 16,
            "priority": "P3", "action": "Connect GSC and baseline the KPIs",
            "owner": "SEO lead", "dependencies": [], "effort": "4h",
            "kpi": "Baseline captured", "approval_class": "agency_admin",
            "status": "Ready", "evidence_ids": [], "confidence": 1.0,
            "implementation_risk": "low", "notes": None,
        },
    ]
    deck = [
        {"kind": "cover", "eyebrow": "ENTERPRISE SEO REVIEW",
         "title": "Evidence first. Growth second.",
         "body": "An approved-domain review of Acme Widgets.", "points": []},
        {"kind": "score", "eyebrow": "EVIDENCE POSTURE",
         "title": "Scores follow coverage.",
         "body": "One category is withheld below the coverage threshold.", "points": []},
        {"kind": "generic", "eyebrow": "TECHNICAL", "title": "Stabilise the crawl surface.",
         "body": "Fix the 404 route and the canonical mismatch first.",
         "points": [{"label": "Errors", "text": "1 broken route"},
                    {"label": "Canonicals", "text": "1 mismatch"}]},
        {"kind": "generic", "eyebrow": "ON-PAGE", "title": "Metadata within budget.",
         "body": "Approve the proposed title rewrites.",
         "points": [{"label": "Titles", "text": "1 oversized"}]},
        {"kind": "generic", "eyebrow": "CONTENT", "title": "Depth where it matters.",
         "body": "Expand the thin about page.", "points": []},
        {"kind": "timeline", "eyebrow": "16-WEEK ROADMAP",
         "title": "Sequence removes risk.",
         "body": "Foundation, hardening, content, measurement."},
        {"kind": "generic", "eyebrow": "MEASUREMENT", "title": "Prove it with baselines.",
         "body": "Connect first-party evidence before forecasting.",
         "points": [{"label": "GSC", "text": "Awaiting connection"}]},
        {"kind": "comparison", "eyebrow": "BEFORE / AFTER",
         "title": "Controls become blockers.",
         "body": "Manual checks become machine-verifiable gates.",
         "points": [{"label": "Before", "text": "Manual spot checks"},
                    {"label": "After", "text": "Release gates with evidence"}]},
        {"kind": "generic", "eyebrow": "NEXT", "title": "Decisions requested.",
         "body": "Approve the plan to unlock implementation.", "points": []},
    ]
    return {
        "schema_version": "1.0.0",
        "client": {"name": "Acme Widgets", "domain": "acmewidgets.com.au",
                   "locale": "en-AU"},
        "project": {"id": "proj-acme", "name": "Acme Widgets Enterprise SEO",
                    "profile": "Enterprise", "business_profile": "ecommerce"},
        "run": {
            "id": "RUN-ACME-001",
            "profile": "Enterprise",
            "configured_page_budget": 500,
            "evidence_as_of": "2026-07-16",
            "captured_at": "2026-07-16T02:00:00+00:00",
            "rule_version": "1.0.0",
            "evidence_coverage": 0.55,
            "coverage_interpretation": "Crawl-only evidence; private sources not connected.",
            "overall_score": None,
            "overall_score_reason": (
                "Weighted evidence coverage 55.0% is below the 70% publication threshold"
            ),
            "state": "GATE_1_REVIEW",
        },
        "executive_summary": (
            "Stabilise the crawl surface, then connect first-party evidence before "
            "expanding content."
        ),
        "sources": [
            {"id": "SRC-CRAWL", "label": "Approved-domain crawl", "kind": "crawl",
             "status": "available", "captured_at": "2026-07-16T02:00:00+00:00",
             "scope": "3 fetched pages", "coverage": 1.0, "unavailable_reason": None},
            {"id": "SRC-GSC", "label": "Google Search Console", "kind": "gsc",
             "status": "unavailable", "captured_at": "2026-07-16T02:00:00+00:00",
             "scope": "Not collected", "coverage": 0.0,
             "unavailable_reason": "credential_not_configured"},
        ],
        "evidence": [],
        "pages": pages,
        "findings": findings,
        "categories": [
            {"category": "Technical", "key": "technical", "score": 72.5,
             "coverage": 1.0, "weight": 0.25, "rule_version": "1.0.0",
             "status": "available", "unavailable_reason": None,
             "evidence_ids": ["EV-0003"]},
            {"category": "Content", "key": "on_page", "score": None,
             "coverage": 0.2, "weight": 0.25, "rule_version": "1.0.0",
             "status": "unavailable",
             "unavailable_reason": "Content evidence coverage below threshold",
             "evidence_ids": []},
        ],
        "content_assets": [
            {"id": "CONTENT-01", "slug": "widget-buying-guide",
             "title": "Widget buying guide", "asset_type": "Editorial guide",
             "target_url": "https://acmewidgets.com.au/shop",
             "audience": "Buyers", "intent": "Commercial",
             "primary_topic": "Widgets", "headline": "Choose the right widget",
             "summary": "A grounded guide.", "body": [], "claims": [],
             "approval_state": "draft_pending_review",
             "generation_method": "grounded", "evidence_ids": ["EV-0003"]},
        ],
        "opportunities": [
            {"id": "OPP-01", "cluster": "Widget comparison",
             "intent": "commercial investigation",
             "target_url": "https://acmewidgets.com.au/shop",
             "decision": "Improve the existing page", "evidence_ids": ["EV-0003"],
             "keyword_volume": None, "ranking": None,
             "unavailable_reason": "GSC and SEMrush not connected"},
            {"id": "OPP-02", "cluster": "Widget care",
             "intent": "informational",
             "target_url": "https://acmewidgets.com.au/about",
             "decision": "Create one supporting article", "evidence_ids": ["EV-0002"],
             "keyword_volume": None, "ranking": None,
             "unavailable_reason": "GSC and SEMrush not connected"},
        ],
        "actions": actions,
        "strategy_sections": [],
        "measurement_plan": [],
        "generation_ledger": [],
        "qa": {"release_status": "PASS", "release_statement": "No blocking failures.",
               "gates": [], "reconciliation": []},
        "limitations": [
            "Private analytics remain unavailable until a connection is approved.",
        ],
        # ---- shared data contract additions -------------------------------
        "market": _market_block(),
        "keywords": [_keyword(index, mapped=index % 2 == 0) for index in range(1, 13)],
        "keyword_clusters": [
            {"id": "CL-01", "name": "Cluster 00", "keyword_count": 6,
             "total_volume": 900, "primary_url": "https://acmewidgets.com.au/shop",
             "intent": "commercial", "coverage": "partial",
             "evidence_ids": ["EV-KW-0001"]},
            {"id": "CL-02", "name": "Cluster 01", "keyword_count": 6,
             "total_volume": 640, "primary_url": None, "intent": "informational",
             "coverage": "gap", "evidence_ids": ["EV-KW-0002"]},
        ],
        "competitors": [
            {"id": f"CMP-{index}", "domain": f"competitor{index}.com.au",
             "relevance": round(0.9 - index * 0.05, 2), "common_keywords": 300 - index * 10,
             "organic_keywords": 5200 - index * 100, "organic_traffic": 21000 - index * 500,
             "organic_cost": 41000.0 - index * 900, "adwords_keywords": 40 + index,
             "gap_keywords": 800 - index * 25, "authority_score": 45 - index,
             "backlinks_total": 30000 - index * 500, "referring_domains": 900 - index * 20,
             "evidence_ids": [f"EV-CMP-{index}"], "unavailable_reason": None}
            for index in range(1, 5)
        ],
        "performance_vs_competitors": {
            "status": "available", "unavailable_reason": None,
            "metrics": [
                {"metric": "Organic keywords", "client": 1840,
                 "competitor_median": 5000, "best_competitor": "competitor1.com.au",
                 "best_value": 5100, "position": "behind",
                 "note": "Provider-reported organic keyword counts."},
                {"metric": "Authority score", "client": 34, "competitor_median": 43,
                 "best_competitor": "competitor1.com.au", "best_value": 44,
                 "position": "behind", "note": "Provider authority score."},
            ],
            "summary": "The client trails the measured competitor median on both metrics.",
        },
        "backlinks": {
            "status": "available", "unavailable_reason": None,
            "overview": _market_block()["domain"],
            "referring_domains": [
                {"domain": f"referrer{index}.com", "authority_score": 20 + index % 60,
                 "backlinks": 10 + index, "country": "AU",
                 "first_seen": "2024-02-01", "last_seen": "2026-07-01"}
                for index in range(1, 9)
            ],
        },
        "onpage_proposals": [
            {"page_id": "URL-0001", "url": "https://acmewidgets.com.au/",
             "page_type": "Homepage",
             "current_title": "Acme Widgets page 1",
             "proposed_title": "Acme Widgets | Australian-made widgets",
             "title_rationale": "Current title is outside the length budget.",
             "current_meta": "Meta description for page 1.",
             "proposed_meta": "A longer, more useful description of the widget range.",
             "meta_rationale": "Current meta description is too short.",
             "current_h1": "Heading 1", "proposed_h1": "Heading 1",
             "h1_rationale": "Observed H1 already states the page subject.",
             "target_keyword": "widget phrase 2", "target_volume": 46,
             "source": "llm_evidence_bound",
             "approval_status": "withheld_pending_editorial_review",
             "evidence_ids": ["EV-0001"]},
            {"page_id": "URL-0002", "url": "https://acmewidgets.com.au/about",
             "page_type": "Content",
             "current_title": "Acme Widgets page 2",
             "proposed_title": "About Acme Widgets | Our Melbourne workshop",
             "title_rationale": "Current title is too short and non-descriptive.",
             "current_meta": None,
             "proposed_meta": "About Acme Widgets and our Melbourne workshop.",
             "meta_rationale": "No meta description was captured.",
             "current_h1": None, "proposed_h1": "About Acme Widgets",
             "h1_rationale": "No H1 was captured on this page.",
             "target_keyword": None, "target_volume": None,
             "source": "deterministic",
             "approval_status": "withheld_pending_editorial_review",
             "evidence_ids": ["EV-0002"]},
        ],
        "crawl_integrity": {
            "status": "clean", "fetched_pages": 3, "challenged_pages": 0,
            "challenge_share": 0.0, "rate_limited_pages": 0, "quarantined_urls": [],
            "note": "No bot challenges were encountered during the crawl window.",
        },
        "deployment": {
            "redirect_candidates": [],
            "canonical_candidates": [
                {"url": "https://acmewidgets.com.au/about",
                 "current_canonical": "https://acmewidgets.com.au/about-us",
                 "proposed_canonical": "https://acmewidgets.com.au/about",
                 "reason": "Canonical points at a variant URL",
                 "approval_status": "review_required", "evidence_id": "EV-0002"},
            ],
            "metadata_review": [
                {"page_id": "URL-0001", "url": "https://acmewidgets.com.au/",
                 "page_type": "Homepage", "status_code": 200,
                 "current_title": "Acme Widgets page 1", "title_length": 76,
                 "title_issue": "Too long",
                 "proposed_title": "Acme Widgets | Australian-made widgets",
                 "current_meta_description": "Meta description for page 1.",
                 "meta_description_length": 28, "meta_description_issue": "Too short",
                 "proposed_meta_description": "A longer, more useful description.",
                 "current_h1": "Heading 1", "h1_issue": "OK",
                 "proposed_h1": "Heading 1",
                 "target_keyword": "Unavailable - GSC and SEMrush not connected",
                 "priority": "P2", "evidence_id": "EV-0001",
                 "approval_status": "withheld_pending_editorial_review"},
            ],
            "internal_link_candidates": [
                {"source_url": "https://acmewidgets.com.au/about",
                 "target_url": "https://acmewidgets.com.au/shop",
                 "anchor": "Browse the widget range",
                 "rationale": "Connect the about narrative to the catalogue",
                 "link_type": "Content to commercial",
                 "observed_status": "Observed in crawl",
                 "evidence_ids": ["EV-0002", "EV-0003"],
                 "approval_status": "review_ready"},
            ],
            "schema": {"deployable": [], "withheld": [
                {"reason": "No verified fact pack was available",
                 "approval_status": "withheld_pending_agency_admin"},
            ]},
            "robots": {"deployable_changes": [],
                       "recommendation": "No robots.txt change is proposed."},
            "disavow": {"enabled": False, "reason": "No backlink evidence connected."},
        },
        "deck": deck,
    }


def _clean_data() -> dict[str, Any]:
    """A no-issue variant: all pages healthy, nothing flagged anywhere."""
    data = _sample_data()
    data["pages"] = [
        _page(1, response_ms=None, body_bytes=None),
        _page(2, response_ms=None, body_bytes=None),
        _page(3, response_ms=None, body_bytes=None,
              normalized_url="https://acmewidgets.com.au/shop-clean",
              canonical_url="https://acmewidgets.com.au/shop-clean"),
    ]
    data["findings"] = []
    data["deployment"]["metadata_review"] = []
    data["deployment"]["canonical_candidates"] = []
    data["deployment"]["internal_link_candidates"] = []
    data["opportunities"] = []
    return data


def _unavailable_data() -> dict[str, Any]:
    """Every optional provider is disconnected; nothing may be invented."""
    data = _sample_data()
    data["market"] = {
        "status": "unavailable", "provider": None, "database": "au",
        "unavailable_reason": "semrush_credential_not_configured",
        "fetched_at": None, "units_spent": 0,
        "domain": dict.fromkeys(_market_block()["domain"]),
    }
    data["keywords"] = []
    data["keyword_clusters"] = []
    data["competitors"] = []
    data["performance_vs_competitors"] = {
        "status": "unavailable",
        "unavailable_reason": "no competitor set was measured for this run",
        "metrics": [], "summary": "Competitor comparison is withheld.",
    }
    data["backlinks"] = {
        "status": "unavailable",
        "unavailable_reason": "no backlink provider is connected",
        "overview": {}, "referring_domains": [],
    }
    data["onpage_proposals"] = []
    data["deployment"]["metadata_review"] = []
    data["crawl_integrity"] = {
        "status": "degraded", "fetched_pages": 3, "challenged_pages": 1,
        "challenge_share": 0.33, "rate_limited_pages": 0,
        "quarantined_urls": ["https://acmewidgets.com.au/blocked"],
        "note": "One page returned a bot challenge and was excluded.",
    }
    return data


def _large_data() -> dict[str, Any]:
    """A synthetic 210-page, 700-keyword dataset used for the V18 volume gates."""
    data = _sample_data()
    pages: list[dict[str, Any]] = []
    for index in range(1, 211):
        if index == 1:
            url = "https://acmewidgets.com.au/"
            page_type = "Homepage"
        elif index <= 121:
            url = f"https://acmewidgets.com.au/products/widget-{index}"
            page_type = "Product"
        elif index <= 171:
            url = f"https://acmewidgets.com.au/collections/range-{index}"
            page_type = "Collection"
        elif index <= 191:
            url = f"https://acmewidgets.com.au/blog/post-{index}"
            page_type = "Editorial"
        else:
            url = f"https://acmewidgets.com.au/pages/info-{index}"
            page_type = "Information"
        page = _page(index)
        page.update({
            "normalized_url": url,
            "original_url": url,
            "original_urls": [url],
            "page_type": page_type,
            "evidence_id": f"EV-{index:04d}",
            "title": f"Acme Widgets {page_type} {index} — Australian made widgets"
            if index % 4 else "Acme Widgets",
            "meta_description": (
                None if index % 5 == 0
                else "A meta description for this page of the Acme Widgets catalogue "
                     "with enough characters to clear the minimum length budget."
            ),
            "h1": None if index % 7 == 0 else f"Heading {index}",
            "word_count": 90 if index % 6 == 0 else 400 + index,
            "canonical_url": (
                f"{url}?variant=1" if index % 9 == 0 else url
            ),
            "status_code": 404 if index % 25 == 0 else 200,
            "schema_types": ["Product"] if page_type == "Product" and index % 3 else [],
            "images_total": 6,
            "images_missing_alt": 2 if index % 3 == 0 else 0,
            "response_ms": 1800 if index % 20 == 0 else 240 + index,
            "body_bytes": 30000 + index * 40,
            "analytics_tags": ["GA4"] if index % 2 else [],
            "external_links": 4,
            "body_sha256": f"hash-{index % 190:03d}",
            "redirect_chain": (
                [f"{url}-old", url] if index % 30 == 0 else [url]
            ),
        })
        pages.append(page)
    data["pages"] = pages

    keywords = [_keyword(index, mapped=index % 3 == 0) for index in range(1, 901)]
    data["keywords"] = keywords
    data["keyword_clusters"] = [
        {"id": f"CL-{index:02d}", "name": f"Cluster {index:02d}",
         "keyword_count": 28, "total_volume": 1200 + index * 40,
         "primary_url": (
             f"https://acmewidgets.com.au/collections/range-{130 + index}"
             if index % 3 else None
         ),
         "intent": ("informational", "commercial", "transactional")[index % 3],
         "coverage": ("covered", "partial", "gap")[index % 3],
         "evidence_ids": [f"EV-CL-{index:02d}"]}
        for index in range(32)
    ]
    data["competitors"] = [
        {"id": f"CMP-{index}", "domain": f"competitor{index}.com.au",
         "relevance": round(0.95 - index * 0.03, 2),
         "common_keywords": 400 - index * 12, "organic_keywords": 6000 - index * 130,
         "organic_traffic": 24000 - index * 400, "organic_cost": 52000.0 - index * 800,
         "adwords_keywords": 30 + index, "gap_keywords": 900 - index * 30,
         "authority_score": 50 - index, "backlinks_total": 40000 - index * 700,
         "referring_domains": 1100 - index * 25,
         "evidence_ids": [f"EV-CMP-{index}"], "unavailable_reason": None}
        for index in range(1, 13)
    ]
    data["backlinks"] = {
        "status": "available", "unavailable_reason": None,
        "overview": _market_block()["domain"],
        "referring_domains": [
            {"domain": f"referrer{index}.example", "authority_score": index % 95,
             "backlinks": 5 + index % 400, "country": "AU",
             "first_seen": "2023-05-02", "last_seen": "2026-07-10"}
            for index in range(1, 1201)
        ],
    }
    data["onpage_proposals"] = [
        {
            "page_id": page["id"], "url": page["normalized_url"],
            "page_type": page["page_type"],
            "current_title": page["title"],
            # Every third page gets a byte-identical "proposal" on purpose: the
            # renderer must refuse to ship it as an optimisation.
            "proposed_title": (
                page["title"] if index % 3 == 0
                else f"{page['page_type']} {index} | Acme Widgets Australia"
            ),
            "title_rationale": "Length and descriptiveness rule.",
            "current_meta": page["meta_description"],
            "proposed_meta": f"Measured description for page {index} of the catalogue.",
            "meta_rationale": "Snippet length rule.",
            "current_h1": page["h1"],
            "proposed_h1": page["h1"] if index % 3 == 0 else f"Acme {page['page_type']} {index}",
            "h1_rationale": "Single descriptive heading rule.",
            "target_keyword": f"widget phrase {index}",
            "target_volume": 40 + index * 3,
            "source": "llm_evidence_bound",
            "approval_status": "withheld_pending_editorial_review",
            "evidence_ids": [page["evidence_id"]],
        }
        for index, page in enumerate(pages, start=1)
    ]
    return data


# --------------------------------------------------------------------------- workbooks


def test_render_workbooks_writes_every_expected_workbook(tmp_path: Path) -> None:
    written = render_workbooks(_sample_data(), tmp_path)
    names = {path.name for path in written}
    assert names >= CORE_WORKBOOKS
    assert ECOMMERCE_WORKBOOK in names  # ecommerce business profile
    assert names >= LINK_WORKBOOKS  # backlink data available
    for path in written:
        assert path.exists(), f"Missing workbook: {path}"
        load_workbook(path)  # must be a valid xlsx

    assert (tmp_path / "01_Audit_Reports" / "Technical_Audit_Report.xlsx").exists()
    assert (tmp_path / "02_Strategy_Documents" / "Master_Keyword_Universe.xlsx").exists()
    assert (tmp_path / "03_Action_Plan" / "16_Week_Action_Plan.xlsx").exists()
    assert (
        tmp_path / "04_Implementation_Deliverables" / "On_Page_Optimizations"
        / "Title_Tag_Optimizations.xlsx"
    ).exists()
    assert (
        tmp_path / "04_Implementation_Deliverables" / "Technical_Fixes"
        / "Redirect_Map.xlsx"
    ).exists()
    assert (tmp_path / "06_QA" / "QC_Report.xlsx").exists()


def test_no_sheet_in_any_workbook_is_empty(tmp_path: Path) -> None:
    for path in render_workbooks(_sample_data(), tmp_path):
        workbook = load_workbook(path)
        for sheet in workbook.worksheets:
            assert _data_row_count(sheet) >= 1, f"{path.name}/{sheet.title} has no rows"
            first = sheet.cell(row=5, column=1).value
            assert first not in (None, ""), f"{path.name}/{sheet.title} row 5 is blank"


def test_workbooks_meet_v18_sheet_and_row_minimums(tmp_path: Path) -> None:
    written = {path.name: path for path in render_workbooks(_large_data(), tmp_path)}
    shortfalls: list[str] = []
    for name, (min_sheets, min_rows) in V18_MINIMUMS.items():
        assert name in written, f"{name} was not rendered"
        workbook = load_workbook(written[name])
        sheets = len(workbook.worksheets)
        rows = sum(_data_row_count(sheet) for sheet in workbook.worksheets)
        if sheets < min_sheets or rows < min_rows:
            shortfalls.append(
                f"{name}: {sheets} sheets / {rows} rows "
                f"(V18 minimum {min_sheets} / {min_rows})"
            )
    assert not shortfalls, "V18 minimums not met:\n" + "\n".join(shortfalls)


def test_local_profile_adds_the_gbp_workbook(tmp_path: Path) -> None:
    data = _sample_data()
    data["project"]["business_profile"] = "hybrid"
    written = {path.name for path in render_workbooks(data, tmp_path)}
    assert LOCAL_WORKBOOK in written
    workbook = load_workbook(tmp_path / "01_Audit_Reports" / LOCAL_WORKBOOK)
    assert len(workbook.worksheets) >= 3
    assert sum(_data_row_count(sheet) for sheet in workbook.worksheets) >= 9
    nap = " ".join(_sheet_texts(workbook["NAP Observations"]))
    assert "GBP API access has not been approved" in nap  # never faked


def test_full_site_inventory_is_one_row_per_page_with_twenty_plus_columns(
    tmp_path: Path,
) -> None:
    data = _large_data()
    render_workbooks(data, tmp_path)
    workbook = load_workbook(tmp_path / "01_Audit_Reports" / "Technical_Audit_Report.xlsx")
    inventory = workbook["Full Site Inventory"]
    headers = [
        inventory.cell(row=4, column=index).value
        for index in range(1, inventory.max_column + 1)
    ]
    assert len([header for header in headers if header]) >= 20
    assert _data_row_count(inventory) == len(data["pages"]) == 210


def test_technical_workbook_structure_brand_band_and_severity_fill(tmp_path: Path) -> None:
    render_workbooks(_sample_data(), tmp_path)
    workbook = load_workbook(
        tmp_path / "01_Audit_Reports" / "Technical_Audit_Report.xlsx"
    )
    assert workbook.sheetnames == [
        "Full Site Inventory", "Error Pages", "Redirects", "Canonical Issues",
        "Duplicate Content", "Indexability", "Title Issues", "Meta Issues",
        "H1 Issues", "Image Alt Issues", "Structured Data", "Findings Register",
        "Methodology",
    ]
    inventory = workbook["Full Site Inventory"]
    assert inventory["A1"].value == "Acme Widgets — Technical Audit Report"
    assert "Run RUN-ACME-001" in str(inventory["A2"].value)
    assert inventory["A4"].value == "URL"
    assert inventory.freeze_panes == "A5"
    assert inventory.auto_filter.ref is not None
    assert inventory["B5"].value == 200  # status codes stay numeric

    errors = workbook["Error Pages"]
    assert errors["A5"].value == "https://acmewidgets.com.au/shop"
    assert errors["B5"].value == 404

    register = workbook["Findings Register"]
    assert register["C5"].value == "Critical"
    assert register["C5"].fill.start_color.rgb == "FFC8443C"

    methodology = workbook["Methodology"]
    texts = _sheet_texts(methodology)
    assert any("1.0.0" in text for text in texts)
    assert any("2026-07-16" in text for text in texts)
    assert any("clean" in text for text in texts)  # crawl integrity is reported


def test_action_plan_exact_columns_and_gantt_span(tmp_path: Path) -> None:
    render_workbooks(_sample_data(), tmp_path)
    workbook = load_workbook(tmp_path / "03_Action_Plan" / "16_Week_Action_Plan.xlsx")
    plan = workbook["Action Plan"]
    headers = [plan.cell(row=4, column=index).value for index in range(1, 14)]
    assert headers == [
        "Phase", "Week", "Task #", "Category", "Description", "Pages/Items",
        "Priority", "Est. Effort", "Owner", "Deliverable", "KPI / Success Metric",
        "Approval Class", "Notes",
    ]
    assert plan.cell(row=5, column=4).value == "Technical"
    assert plan.cell(row=9, column=4).value == "General"
    assert plan.cell(row=5, column=7).fill.start_color.rgb == "FFC8443C"

    gantt = workbook["Gantt"]
    assert gantt.cell(row=7, column=5).fill.start_color.rgb == "FF1B7CA8"
    assert gantt.cell(row=7, column=9).fill.start_color.rgb != "FF1B7CA8"


def test_empty_sections_render_honest_statements(tmp_path: Path) -> None:
    render_workbooks(_clean_data(), tmp_path)
    workbook = load_workbook(
        tmp_path / "01_Audit_Reports" / "Technical_Audit_Report.xlsx"
    )
    assert workbook["Error Pages"]["A5"].value == EMPTY_MESSAGE
    assert workbook["Redirects"]["A5"].value == EMPTY_MESSAGE
    assert workbook["Duplicate Content"]["A5"].value == EMPTY_MESSAGE

    baseline = load_workbook(
        tmp_path / "01_Audit_Reports" / "Baseline_Performance_Analysis.xlsx"
    )
    assert baseline["Response Times"]["A5"].value == (
        "Unavailable — response timing was not captured for this run."
    )

    content = load_workbook(
        tmp_path / "01_Audit_Reports" / "Content_Audit_Workbook.xlsx"
    )
    assert content["Duplicate or Thin Content"]["A5"].value == EMPTY_MESSAGE


def test_formula_injection_is_neutralised(tmp_path: Path) -> None:
    data = _sample_data()
    data["pages"][0]["title"] = "=HYPERLINK('http://evil.example','click')"
    render_workbooks(data, tmp_path)
    workbook = load_workbook(
        tmp_path / "01_Audit_Reports" / "Technical_Audit_Report.xlsx"
    )
    cell = workbook["Full Site Inventory"]["C5"]
    assert cell.data_type != "f"
    assert str(cell.value).startswith("'=")


def test_keyword_universe_renders_every_measured_keyword(tmp_path: Path) -> None:
    data = _large_data()
    render_workbooks(data, tmp_path)
    workbook = load_workbook(
        tmp_path / "02_Strategy_Documents" / "Master_Keyword_Universe.xlsx"
    )
    mapping = workbook["Keyword Research Mapping"]
    assert _data_row_count(mapping) == len(data["keywords"]) == 900
    assert mapping["B5"].value == "widget phrase 1"
    assert mapping["E5"].value == 43  # measured volume stays numeric
    funnel = workbook["Funnel Distribution"]
    assert {funnel.cell(row=row, column=1).value for row in range(5, 8)} == {
        "TOFU", "MOFU", "BOFU",
    }


def test_unavailable_providers_render_reasons_not_estimates(tmp_path: Path) -> None:
    data = _unavailable_data()
    written = {path.name for path in render_workbooks(data, tmp_path)}
    assert written >= CORE_WORKBOOKS
    assert LINK_WORKBOOKS.isdisjoint(written)  # no link building without link data

    backlinks = load_workbook(
        tmp_path / "01_Audit_Reports" / "Backlink_Audit_Report.xlsx"
    )
    assert {"Overview", "Referring Domains", "Authority Distribution",
            "Link Gap Opportunities", "Competitor Comparison",
            "Methodology"} <= set(backlinks.sheetnames)
    for name in ("Overview", "Referring Domains", "Authority Distribution"):
        assert backlinks[name]["A5"].value == (
            "Unavailable — no backlink provider is connected"
        )
    methodology_text = " ".join(_sheet_texts(backlinks["Methodology"]))
    assert "connected provider adds" in methodology_text

    keywords = load_workbook(
        tmp_path / "02_Strategy_Documents" / "Master_Keyword_Universe.xlsx"
    )
    assert keywords["Keyword Research Mapping"]["A5"].value == (
        "Unavailable — semrush_credential_not_configured"
    )
    competitors = load_workbook(
        tmp_path / "01_Audit_Reports" / "Competitor_Landscape_Analysis.xlsx"
    )
    assert competitors["Performance vs Competitors"]["A5"].value == (
        "Unavailable — no competitor set was measured for this run"
    )


def test_identical_proposals_are_marked_no_change_required(tmp_path: Path) -> None:
    render_workbooks(_large_data(), tmp_path)
    onpage = tmp_path / "04_Implementation_Deliverables" / "On_Page_Optimizations"
    for filename, sheet_name in (
        ("Title_Tag_Optimizations.xlsx", "Title Tags"),
        ("H1_Tags.xlsx", "H1 Tags"),
    ):
        workbook = load_workbook(onpage / filename)
        sheet = workbook[sheet_name]
        statuses: list[str] = []
        for row in range(5, 5 + _data_row_count(sheet)):
            current = sheet.cell(row=row, column=4).value
            proposed = sheet.cell(row=row, column=6).value
            status = sheet.cell(row=row, column=8).value
            statuses.append(str(status))
            if status == NO_CHANGE:
                assert proposed == NO_CHANGE
            elif status == "Proposed":
                assert str(proposed).strip().casefold() != str(current).strip().casefold()
        assert NO_CHANGE in statuses, f"{filename} never exercises the no-change path"
        assert "Proposed" in statuses


def test_deployment_fix_sheets_are_one_to_one_with_observations(tmp_path: Path) -> None:
    data = _large_data()
    render_workbooks(data, tmp_path)
    fixes = tmp_path / "04_Implementation_Deliverables" / "Technical_Fixes"

    errors = [
        page for page in data["pages"]
        if isinstance(page["status_code"], int) and page["status_code"] >= 400
    ]
    redirects = load_workbook(fixes / "Redirect_Map.xlsx")["Redirect Map"]
    assert _data_row_count(redirects) == len(errors) > 0

    mismatched = [
        page for page in data["pages"]
        if str(page["canonical_url"]).rstrip("/") != str(page["normalized_url"]).rstrip("/")
    ]
    canonical = load_workbook(fixes / "Canonical_Fixes.xlsx")["Canonical Fixes"]
    assert _data_row_count(canonical) >= len(mismatched) > 0


def test_geo_scores_are_rule_derived_and_never_blank(tmp_path: Path) -> None:
    render_workbooks(_large_data(), tmp_path)
    workbook = load_workbook(
        tmp_path / "01_Audit_Reports" / "GEO_AEO_Readiness_Scorecard.xlsx"
    )
    readiness = workbook["Page-Level Readiness"]
    for row in range(5, 5 + _data_row_count(readiness)):
        score = readiness.cell(row=row, column=8).value
        assert isinstance(score, int)
        assert 0 <= score <= 100
        assert readiness.cell(row=row, column=9).value  # score basis is always stated


# --------------------------------------------------------------------------- pptx deck


def test_deck_renders_one_slide_per_spec_with_footers(tmp_path: Path) -> None:
    data = _sample_data()
    output = render_deck(data, tmp_path / "07_Executive_Deck" / "Executive_Deck.pptx")
    assert output.exists()
    presentation = Presentation(str(output))
    assert len(presentation.slides) == len(data["deck"]) == 9
    assert presentation.slide_width == Inches(13.333)
    assert presentation.slide_height == Inches(7.5)

    combined = "\n".join(_deck_texts(presentation))
    assert "Evidence first. Growth second." in combined
    assert "slide 1/9" in combined
    assert "slide 9/9" in combined
    assert "Acme Widgets · Enterprise SEO Audit" in combined
    assert "None" not in combined


def test_deck_never_repeats_a_layout_on_consecutive_slides() -> None:
    layouts = slide_layouts_for(_sample_data())
    assert len(layouts) == 9
    assert len(set(layouts)) >= 5
    repeats = [
        (index, layouts[index])
        for index in range(1, len(layouts))
        if layouts[index] == layouts[index - 1]
    ]
    assert not repeats, f"consecutive layout repeats: {repeats}"

    # A deck that asks for the same semantic kind twice in a row still varies.
    data = _sample_data()
    data["deck"] = [dict(data["deck"][1]) for _ in range(4)]
    stressed = slide_layouts_for(data)
    assert all(
        stressed[index] != stressed[index - 1] for index in range(1, len(stressed))
    )


def test_deck_uses_native_charts(tmp_path: Path) -> None:
    output = render_deck(_sample_data(), tmp_path / "deck.pptx")
    presentation = Presentation(str(output))
    chart_types = []
    for slide in presentation.slides:
        for shape in slide.shapes:
            if shape.has_chart:
                chart_types.append(shape.chart.chart_type)
    assert len(chart_types) >= 3, "expected native bar, doughnut and comparison charts"
    names = {str(item) for item in chart_types}
    assert any("BAR" in name for name in names)
    assert any("DOUGHNUT" in name for name in names)
    assert any("COLUMN" in name for name in names)


def test_deck_score_slide_renders_withheld_category_with_reason(tmp_path: Path) -> None:
    data = _sample_data()
    output = render_deck(data, tmp_path / "deck.pptx")
    presentation = Presentation(str(output))
    score_slide = presentation.slides[1]
    slide_text = "\n".join(
        shape.text_frame.text for shape in score_slide.shapes if shape.has_text_frame
    )
    assert "Withheld" in slide_text
    assert "Content evidence coverage below threshold" in slide_text

    # The withheld category must never reach the chart as a zero.
    charted: list[float] = []
    for shape in score_slide.shapes:
        if shape.has_chart:
            for series in shape.chart.plots[0].series:
                charted.extend(float(value) for value in series.values)
    assert charted == [72.5]
    assert 0.0 not in charted


def test_deck_timeline_slide_is_a_sixteen_week_gantt(tmp_path: Path) -> None:
    output = render_deck(_sample_data(), tmp_path / "deck.pptx")
    presentation = Presentation(str(output))
    timeline_slide = presentation.slides[5]
    slide_text = "\n".join(
        shape.text_frame.text for shape in timeline_slide.shapes if shape.has_text_frame
    )
    assert "Foundation" in slide_text
    assert "Measurement" in slide_text
    assert "Weeks 13–16" in slide_text
    for week in ("W1", "W8", "W16"):
        assert week in slide_text


def test_deck_theme_fonts_are_branded(tmp_path: Path) -> None:
    import zipfile

    output = render_deck(_sample_data(), tmp_path / "deck.pptx")
    with zipfile.ZipFile(output) as archive:
        theme = archive.read("ppt/theme/theme1.xml").decode("utf-8")
    assert "Fraunces" in theme
    assert "Source Sans 3" in theme
    assert '<a:majorFont><a:latin typeface="Calibri"' not in theme


# --------------------------------------------------------------------------- deck pdf


def test_deck_pdf_has_one_page_per_slide(tmp_path: Path) -> None:
    data = _sample_data()
    output = render_deck_pdf(data, tmp_path / "07_Executive_Deck" / "Executive_Deck.pdf")
    assert output.exists()
    reader = PdfReader(str(output))
    assert len(reader.pages) == len(data["deck"]) == 9
    first_page = reader.pages[0].extract_text() or ""
    assert "Evidence first" in first_page


def test_deck_pdf_withholds_scores_rather_than_zeroing_them(tmp_path: Path) -> None:
    output = render_deck_pdf(_sample_data(), tmp_path / "deck.pdf")
    reader = PdfReader(str(output))
    text = "\n".join(page.extract_text() or "" for page in reader.pages)
    assert "Withheld" in text
    assert "Content evidence coverage below threshold" in text
    assert "slide 9/9" in text


def test_deck_pdf_is_honest_when_providers_are_disconnected(tmp_path: Path) -> None:
    output = render_deck_pdf(_unavailable_data(), tmp_path / "deck.pdf")
    text = "\n".join(page.extract_text() or "" for page in PdfReader(str(output)).pages)
    assert "no competitor set was measured for this run" in text


# --------------------------------------------------------------------------- markdown


def test_markdown_contains_required_sections_and_no_none_literal() -> None:
    markdown = render_markdown(_sample_data())
    assert markdown.startswith("# Acme Widgets — Enterprise SEO Audit Results")
    for header in [
        "## At a glance",
        "## Category scorecard",
        "## Top priority findings",
        "## 16-week action plan overview",
        "## What is in this package",
        "## Data sources & coverage",
        "## Methodology & limitations",
    ]:
        assert header in markdown
    assert "None" not in markdown
    assert "\r" not in markdown
    assert markdown.endswith("\n")
    assert "Generated by Traffic Radius Enterprise SEO Studio · run RUN-ACME-001" in markdown
    assert "| Metric | Value |" in markdown
    assert "| Category | Score | Evidence coverage | Findings |" in markdown
    assert "```" in markdown
    assert "Technical_Audit_Report.xlsx" in markdown
    assert "Schema_Product_Template.json" in markdown
    assert "Schema_LocalBusiness.json" not in markdown  # ecommerce profile only
    assert "CONTENT-01_widget-buying-guide.docx" in markdown


def test_markdown_renders_withheld_scores_with_reason() -> None:
    markdown = render_markdown(_sample_data())
    assert (
        "Withheld — Weighted evidence coverage 55.0% is below the 70% "
        "publication threshold" in markdown
    )
    scorecard_row = next(
        line for line in markdown.splitlines() if line.startswith("| Content |")
    )
    cells = [cell.strip() for cell in scorecard_row.strip("|").split("|")]
    assert cells[1] == "Withheld"


def test_markdown_honest_when_run_is_empty() -> None:
    data = _clean_data()
    data["actions"] = []
    data["categories"] = []
    data["content_assets"] = []
    markdown = render_markdown(data)
    assert "No actions were scheduled for this run." in markdown
    assert "No categories were scored in this run" in markdown
    assert "no content assets cleared evidence checks" in markdown
    assert "None" not in markdown


# ---------------------------------------------------------------------------
# CRO signal quality: normal ecommerce copy is not a defect; floods aggregate
# ---------------------------------------------------------------------------


def _cro_page(word_count: int, page_type: str = "Product", index: int = 0) -> dict:
    return {
        "normalized_url": f"https://example.com.au/p/{index}",
        "page_type": page_type,
        "word_count": word_count,
        "response_ms": 300,
        "images_missing_alt": 0,
        "images_total": 4,
        "meta_description": "Present",
        "h1": "Heading",
        "facts": {"has_viewport": True, "h1_count": 1},
        "evidence_id": f"EV-{index:04d}",
    }


def test_normal_product_copy_is_not_flagged_thin() -> None:
    from exporters.xlsx_workbooks import _cro_findings

    findings = _cro_findings([_cro_page(200), _cro_page(180, index=1)])
    assert not [row for row in findings if row[2] == "Thin money page"]


def test_truly_thin_product_page_is_flagged_medium_then_high() -> None:
    from exporters.xlsx_workbooks import _cro_findings

    medium = _cro_findings([_cro_page(110)])
    assert [row[3] for row in medium if row[2] == "Thin money page"] == ["Medium"]

    high = _cro_findings([_cro_page(40)])
    assert [row[3] for row in high if row[2] == "Thin money page"] == ["High"]


def test_signal_floods_collapse_to_one_summary_row() -> None:
    from exporters.xlsx_workbooks import _aggregate_signal_floods, _cro_findings

    pages = [_cro_page(90, index=i) for i in range(22)]
    aggregated = _aggregate_signal_floods(_cro_findings(pages))
    thin_rows = [row for row in aggregated if row[2] == "Thin money page"]

    assert len(thin_rows) == 1
    assert "22 pages" in str(thin_rows[0][0])
    assert "template-level fix" in str(thin_rows[0][4])
