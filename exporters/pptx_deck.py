"""Executive deck renderer (python-pptx, 16:9) plus its native PDF sibling.

Pure ``dict -> file``: consumes the compiled run-data dictionary and renders
``data['deck']`` slide specs onto blank layouts with explicit shapes.

Quality rules enforced here:

* The presentation theme's major/minor fonts are patched to the brand pair, so
  the deck does not silently ship in Calibri, and every run carries an explicit
  font name as well.
* Charts are **native** OOXML charts (``CategoryChartData`` +
  ``XL_CHART_TYPE``), not drawn rectangles: a bar chart for category scores, a
  doughnut for the severity mix and a clustered bar for competitor performance.
* Withheld scores are never coerced to zero. They are excluded from chart data
  and rendered as a "Withheld" chip with the withholding reason.
* Layout variety is enforced: the same layout never appears on two consecutive
  slides.
* Every slide carries the brand mark (drawn natively from the three bars of
  ``static/img/traffic-radius-mark.svg``) and a footer with client, evidence
  date and slide n/N.
"""

from __future__ import annotations

import re
from pathlib import Path
from typing import Any

from pptx import Presentation
from pptx.chart.data import CategoryChartData
from pptx.dml.color import RGBColor
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.opc.constants import RELATIONSHIP_TYPE as RT
from pptx.util import Emu, Inches, Pt

from exporters.brand import (
    ACCENT_LIGHT,
    BRAND_BLUE,
    BRAND_DEEP,
    BRAND_GREEN,
    CRITICAL,
    GOLD,
    GOLD_LIGHT,
    MUTED_ON_NAVY,
    NAVY_500,
    NAVY_700,
    NAVY_900,
    NEUTRAL,
    POSITIVE,
    RULE_ON_NAVY,
    TEXT_ON_NAVY,
    WARNING,
    WHITE,
    BrandFonts,
)
from exporters.common import safe_text

SLIDE_WIDTH_IN = 13.333
SLIDE_HEIGHT_IN = 7.5
MARGIN_IN = 0.6
CONTENT_WIDTH_IN = SLIDE_WIDTH_IN - 2 * MARGIN_IN

_FONTS = BrandFonts()

#: Layouts a generic slide rotates through. Consecutive repeats are forbidden.
GENERIC_LAYOUTS: tuple[str, ...] = ("statement", "stat_rail", "chart", "two_column", "table")

#: Fallback layout when a semantic layout would repeat back to back.
LAYOUT_FALLBACK: dict[str, str] = {
    "cover": "statement",
    "score": "stat_rail",
    "timeline": "table",
    "comparison": "two_column",
}

SEVERITY_ORDER: tuple[str, ...] = ("Critical", "High", "Medium", "Low")
SEVERITY_CHART_COLORS: dict[str, str] = {
    "Critical": CRITICAL,
    "High": WARNING,
    "Medium": ACCENT_LIGHT,
    "Low": NEUTRAL[300],
}


def _rgb(color: str) -> RGBColor:
    return RGBColor.from_string(color.lstrip("#"))


# --------------------------------------------------------------------------- primitives


def _background(slide: Any, color: str) -> None:
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = _rgb(color)


def _textbox(
    slide: Any,
    left: float,
    top: float,
    width: float,
    height: float,
    text: str,
    *,
    size: int,
    color: str,
    bold: bool = False,
    italic: bool = False,
    font: str | None = None,
    align: PP_ALIGN = PP_ALIGN.LEFT,
) -> Any:
    box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    frame = box.text_frame
    frame.word_wrap = True
    paragraph = frame.paragraphs[0]
    paragraph.alignment = align
    run = paragraph.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = _rgb(color)
    run.font.name = font or _FONTS.body
    return box


def _rect(
    slide: Any,
    left: float,
    top: float,
    width: float,
    height: float,
    fill_color: str | None,
    *,
    line_color: str | None = None,
    line_width: float = 0.75,
    shape: MSO_SHAPE = MSO_SHAPE.RECTANGLE,
) -> Any:
    element = slide.shapes.add_shape(
        shape, Inches(left), Inches(top), Inches(width), Inches(height)
    )
    element.shadow.inherit = False
    if fill_color is None:
        element.fill.background()
    else:
        element.fill.solid()
        element.fill.fore_color.rgb = _rgb(fill_color)
    if line_color is None:
        element.line.fill.background()
    else:
        element.line.color.rgb = _rgb(line_color)
        element.line.width = Pt(line_width)
    return element


def _set_fill_opacity(shape: Any, opacity: float) -> None:
    """Apply an ``a:alpha`` child to a solid fill (0.0 transparent … 1.0 opaque)."""
    from pptx.oxml.ns import qn

    solid = shape.fill._xPr.find(qn("a:solidFill"))
    if solid is None:  # pragma: no cover - defensive
        return
    srgb = solid.find(qn("a:srgbClr"))
    if srgb is None:  # pragma: no cover - defensive
        return
    value = str(int(max(0.0, min(1.0, opacity)) * 100000))
    srgb.append(srgb.makeelement(qn("a:alpha"), {"val": value}))


def _dark_chart(chart: Any) -> None:
    """Make the chart plate transparent so the navy slide shows through."""
    from lxml import etree
    from pptx.oxml.ns import qn

    chart_space = chart._chartSpace
    if chart_space.find(qn("c:spPr")) is not None:  # pragma: no cover - defensive
        return
    sp_pr = etree.SubElement(chart_space, qn("c:spPr"))
    etree.SubElement(sp_pr, qn("a:noFill"))
    line = etree.SubElement(sp_pr, qn("a:ln"))
    etree.SubElement(line, qn("a:noFill"))
    chart_space.find(qn("c:chart")).addnext(sp_pr)


def _chip(slide: Any, left: float, top: float, width: float, height: float,
          text: str, *, fill_color: str, text_color: str, size: int = 11) -> Any:
    shape = _rect(slide, left, top, width, height, fill_color,
                  shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    frame = shape.text_frame
    frame.word_wrap = False
    run = frame.paragraphs[0].add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = True
    run.font.name = _FONTS.body
    run.font.color.rgb = _rgb(text_color)
    return shape


def _brand_mark(slide: Any, *, dark: bool) -> None:
    """Draw the three bars of the Traffic Radius mark as native shapes."""
    left = SLIDE_WIDTH_IN - MARGIN_IN - 0.62
    base = 0.86
    bars = (
        (0.0, 0.22, WHITE if dark else BRAND_DEEP),
        (0.22, 0.34, BRAND_BLUE),
        (0.44, 0.46, BRAND_GREEN),
    )
    for offset, height, color in bars:
        _rect(slide, left + offset, base - height, 0.15, height, color,
              shape=MSO_SHAPE.ROUNDED_RECTANGLE)


def _footer(slide: Any, data: dict[str, Any], index: int, total: int, *, dark: bool) -> None:
    client = safe_text(data.get("client", {}).get("name"), "Client")
    as_of = safe_text(data.get("run", {}).get("evidence_as_of"))
    text = (
        f"{client} · Enterprise SEO Audit · Evidence as of {as_of} · "
        f"slide {index}/{total}"
    )
    del dark  # every slide now sits on navy, so the footer tone is constant
    _textbox(
        slide, MARGIN_IN, SLIDE_HEIGHT_IN - 0.42, CONTENT_WIDTH_IN, 0.3, text,
        size=9, color=MUTED_ON_NAVY,
    )


def _eyebrow(slide: Any, text: str, *, top: float = 0.55) -> None:
    if not text:
        return
    _textbox(
        slide, MARGIN_IN, top, CONTENT_WIDTH_IN, 0.35, text.upper(),
        size=13, color=GOLD, bold=True,
    )


#: Titles longer than this are assumed to wrap onto a second line at 30pt.
_TITLE_WRAP_CHARS = 45


def _title(slide: Any, text: str, *, top: float = 0.95, color: str = TEXT_ON_NAVY,
           size: int = 30) -> float:
    """Draw the slide title and return the recommended body top (inches).

    Long titles wrap to a second line, so every renderer that anchors body
    content below the title band must start at the returned offset instead of a
    hard-coded 2.0in.
    """
    _textbox(
        slide, MARGIN_IN, top, CONTENT_WIDTH_IN, 1.3, text,
        size=size, color=color, bold=True, font=_FONTS.display,
    )
    return 2.35 if len(text) > _TITLE_WRAP_CHARS else 2.0


# --------------------------------------------------------------------------- data views


def _scored_categories(data: dict[str, Any]) -> tuple[list[dict[str, Any]],
                                                      list[dict[str, Any]]]:
    """Split categories into (scored, withheld). A null score is never zeroed."""
    scored: list[dict[str, Any]] = []
    withheld: list[dict[str, Any]] = []
    for category in list(data.get("categories") or [])[:8]:
        if isinstance(category.get("score"), int | float):
            scored.append(category)
        else:
            withheld.append(category)
    return scored, withheld


def _severity_mix(data: dict[str, Any]) -> list[tuple[str, int]]:
    counts: dict[str, int] = {}
    for finding in data.get("findings") or []:
        label = str(finding.get("severity") or "").strip().title()
        if not label:
            continue
        counts[label] = counts.get(label, 0) + 1
    ordered = [(name, counts[name]) for name in SEVERITY_ORDER if name in counts]
    extras = sorted(
        (item for item in counts.items() if item[0] not in SEVERITY_ORDER),
        key=lambda item: -item[1],
    )
    return ordered + extras


def _withheld_reason(data: dict[str, Any], withheld: list[dict[str, Any]]) -> str | None:
    if not withheld:
        return None  # nothing was withheld, so no footnote is owed
    for category in withheld:
        reason = category.get("unavailable_reason")
        if reason:
            return str(reason)
    reason = data.get("run", {}).get("overall_score_reason")
    return str(reason) if reason else None


# --------------------------------------------------------------------------- layouts


def _resolve_layouts(specs: list[dict[str, Any]]) -> list[str]:
    """Assign one layout per slide, never repeating a layout back to back."""
    layouts: list[str] = []
    previous: str | None = None
    rotation = 0
    for spec in specs:
        kind = str(spec.get("kind") or "generic").casefold()
        if kind in {"cover", "score", "timeline", "comparison"}:
            layout = kind
        else:
            layout = GENERIC_LAYOUTS[rotation % len(GENERIC_LAYOUTS)]
            rotation += 1
        if layout == previous:
            layout = LAYOUT_FALLBACK.get(layout, "")
            if not layout or layout == previous:
                layout = next(
                    candidate for candidate in GENERIC_LAYOUTS if candidate != previous
                )
                rotation += 1
        layouts.append(layout)
        previous = layout
    return layouts


# --------------------------------------------------------------------------- renderers


def _render_cover(slide: Any, spec: dict[str, Any], data: dict[str, Any]) -> None:
    _background(slide, NAVY_900)
    client = safe_text(data.get("client", {}).get("name"), "Client")
    as_of = safe_text(data.get("run", {}).get("evidence_as_of"))
    # Depth object first so every later shape stacks above it.
    oval = _rect(slide, 9.7, 4.3, 5.6, 5.6, NAVY_500, shape=MSO_SHAPE.OVAL)
    _set_fill_opacity(oval, 0.16)
    _rect(slide, 0, 0, 0.16, SLIDE_HEIGHT_IN, GOLD)
    _textbox(
        slide, MARGIN_IN, 1.7, CONTENT_WIDTH_IN, 0.4,
        safe_text(spec.get("eyebrow"), "ENTERPRISE SEO REVIEW").upper(),
        size=14, color=GOLD, bold=True,
    )
    _textbox(
        slide, MARGIN_IN, 2.15, CONTENT_WIDTH_IN, 2.2,
        safe_text(spec.get("title"), "Enterprise SEO Audit"),
        size=48, color=TEXT_ON_NAVY, bold=True, font=_FONTS.display,
    )
    _rect(slide, MARGIN_IN, 4.55, 2.2, 0.06, GOLD)
    body = spec.get("body")
    if body:
        _textbox(slide, MARGIN_IN, 4.75, CONTENT_WIDTH_IN - 2.0, 1.0, str(body),
                 size=18, color=TEXT_ON_NAVY)
    _rect(slide, MARGIN_IN, 5.78, CONTENT_WIDTH_IN, 0.015, RULE_ON_NAVY)
    _textbox(
        slide, MARGIN_IN, 5.9, CONTENT_WIDTH_IN, 0.4,
        f"{client} · Evidence as of {as_of}", size=14, color=MUTED_ON_NAVY,
    )


def _render_score(slide: Any, spec: dict[str, Any], data: dict[str, Any],
                  body_top: float) -> None:
    """Native bar chart of scored categories; withheld ones stay text chips."""
    body = spec.get("body")
    if body:
        _textbox(slide, MARGIN_IN, body_top, CONTENT_WIDTH_IN, 0.55, str(body),
                 size=16, color=TEXT_ON_NAVY)
    scored, withheld = _scored_categories(data)
    chart_top = body_top + 0.65
    # With nothing withheld the right panel is not reserved: full-width chart.
    chart_width = CONTENT_WIDTH_IN * 0.62 if withheld else CONTENT_WIDTH_IN

    if scored:
        chart_data = CategoryChartData()
        chart_data.categories = [safe_text(item.get("category")) for item in scored]
        chart_data.add_series(
            "Score", tuple(float(item["score"]) for item in scored)
        )
        graphic_frame = slide.shapes.add_chart(
            XL_CHART_TYPE.BAR_CLUSTERED,
            Inches(MARGIN_IN), Inches(chart_top),
            Inches(chart_width), Inches(3.5),
            chart_data,
        )
        chart = graphic_frame.chart
        _dark_chart(chart)
        chart.has_legend = False
        chart.has_title = False
        plot = chart.plots[0]
        plot.gap_width = 60
        plot.vary_by_categories = False
        series = plot.series[0]
        series.format.fill.solid()
        series.format.fill.fore_color.rgb = _rgb(GOLD)
        value_axis = chart.value_axis
        value_axis.minimum_scale = 0
        value_axis.maximum_scale = 100
        for axis in (chart.category_axis, chart.value_axis):
            axis.tick_labels.font.size = Pt(11)
            axis.tick_labels.font.name = _FONTS.body
            axis.tick_labels.font.color.rgb = _rgb(MUTED_ON_NAVY)
    else:
        _textbox(
            slide, MARGIN_IN, body_top + 0.9, chart_width, 0.5,
            "No category cleared the evidence threshold, so no score is published.",
            size=16, color=MUTED_ON_NAVY,
        )

    if withheld:
        panel_left = MARGIN_IN + CONTENT_WIDTH_IN * 0.66
        panel_width = CONTENT_WIDTH_IN - CONTENT_WIDTH_IN * 0.66
        _textbox(slide, panel_left, chart_top, panel_width, 0.32, "WITHHELD",
                 size=12, color=MUTED_ON_NAVY, bold=True)
        top = chart_top + 0.4
        for category in withheld:
            _textbox(slide, panel_left, top, panel_width, 0.3,
                     safe_text(category.get("category")), size=13,
                     color=TEXT_ON_NAVY, bold=True)
            _chip(slide, panel_left, top + 0.32, 1.35, 0.3, "Withheld",
                  fill_color=NAVY_500, text_color=TEXT_ON_NAVY)
            top += 0.85

    reason = _withheld_reason(data, withheld)
    if reason:  # _withheld_reason is None when nothing is withheld
        _textbox(
            slide, MARGIN_IN, chart_top + 3.6, CONTENT_WIDTH_IN, 0.5,
            f"Withheld: {reason}", size=11, color=MUTED_ON_NAVY, italic=True,
        )


def _render_stat_rail(slide: Any, spec: dict[str, Any], data: dict[str, Any],
                      body_top: float) -> None:
    """A rail of stat cards: scored categories first, then the slide's own points."""
    body = spec.get("body")
    if body:
        _textbox(slide, MARGIN_IN, body_top, CONTENT_WIDTH_IN, 0.6, str(body),
                 size=16, color=TEXT_ON_NAVY)
    scored, withheld = _scored_categories(data)
    cards: list[tuple[str, str, str]] = [
        (safe_text(item.get("category")), f"{float(item['score']):.0f}", "score")
        for item in scored[:4]
    ]
    for item in withheld:
        if len(cards) >= 4:
            break
        cards.append((safe_text(item.get("category")), "Withheld", "withheld"))
    for point in list(spec.get("points") or []):
        if len(cards) >= 4:
            break
        cards.append((safe_text(point.get("label")), safe_text(point.get("text")), "point"))
    if not cards:
        _textbox(slide, MARGIN_IN, body_top + 1.0, CONTENT_WIDTH_IN, 0.5,
                 "No measured statistics were available for this slide.",
                 size=15, color=MUTED_ON_NAVY, italic=True)
        return

    gap = 0.3
    card_top = body_top + 0.9
    card_height = 3.4
    width = (CONTENT_WIDTH_IN - gap * (len(cards) - 1)) / len(cards)
    for index, (label, value, kind) in enumerate(cards):
        left = MARGIN_IN + index * (width + gap)
        _rect(slide, left, card_top, width, card_height, NAVY_700,
              line_color=RULE_ON_NAVY)
        # One gold strip on every card: the number carries the signal.
        _rect(slide, left, card_top, width, 0.09, GOLD)
        _textbox(slide, left + 0.22, card_top + 0.25, width - 0.44, 0.4, label,
                 size=13, color=MUTED_ON_NAVY, bold=True)
        value_box = _textbox(
            slide, left + 0.22, card_top + 0.7, width - 0.44, card_height - 0.95,
            value,
            size=30 if kind == "score" else 15,
            color=GOLD_LIGHT if kind == "score"
            else (MUTED_ON_NAVY if kind == "withheld" else TEXT_ON_NAVY),
            bold=kind == "score", font=_FONTS.display if kind == "score" else None,
        )
        value_box.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
    reason = _withheld_reason(data, withheld)
    if reason and any(kind == "withheld" for _, _, kind in cards):
        _textbox(slide, MARGIN_IN, card_top + card_height + 0.12, CONTENT_WIDTH_IN, 0.4,
                 f"Withheld: {reason}", size=11, color=MUTED_ON_NAVY, italic=True)


def _render_chart(slide: Any, spec: dict[str, Any], data: dict[str, Any],
                  body_top: float) -> None:
    """Native doughnut of the severity mix, with the slide's points beside it."""
    body = spec.get("body")
    if body:
        _textbox(slide, MARGIN_IN, body_top, CONTENT_WIDTH_IN, 0.6, str(body),
                 size=16, color=TEXT_ON_NAVY)
    mix = _severity_mix(data)
    if mix:
        chart_data = CategoryChartData()
        chart_data.categories = [name for name, _ in mix]
        chart_data.add_series("Findings", tuple(count for _, count in mix))
        graphic_frame = slide.shapes.add_chart(
            XL_CHART_TYPE.DOUGHNUT,
            Inches(MARGIN_IN), Inches(body_top + 0.7),
            Inches(5.0), Inches(3.6), chart_data,
        )
        chart = graphic_frame.chart
        _dark_chart(chart)
        chart.has_title = False
        chart.has_legend = True
        chart.legend.position = XL_LEGEND_POSITION.RIGHT
        chart.legend.include_in_layout = False
        chart.legend.font.size = Pt(12)
        chart.legend.font.name = _FONTS.body
        chart.legend.font.color.rgb = _rgb(TEXT_ON_NAVY)
        points = chart.plots[0].series[0].points
        for index, (name, _count) in enumerate(mix):
            point = points[index]
            point.format.fill.solid()
            point.format.fill.fore_color.rgb = _rgb(
                SEVERITY_CHART_COLORS.get(name, GOLD if index % 2 else GOLD_LIGHT)
            )
    else:
        _textbox(slide, MARGIN_IN, body_top + 1.0, 5.0, 0.6,
                 "No findings were raised, so there is no severity mix to chart.",
                 size=15, color=MUTED_ON_NAVY, italic=True)

    left = MARGIN_IN + 5.6
    width = CONTENT_WIDTH_IN - 5.6
    top = body_top + 0.75
    for name, count in mix[:4]:
        _rect(slide, left, top, width, 0.7, NAVY_700, line_color=RULE_ON_NAVY)
        _rect(slide, left, top, 0.09, 0.7,
              SEVERITY_CHART_COLORS.get(name, GOLD))
        _textbox(slide, left + 0.25, top + 0.1, width - 0.5, 0.3, name,
                 size=13, color=TEXT_ON_NAVY, bold=True)
        _textbox(slide, left + 0.25, top + 0.38, width - 0.5, 0.28,
                 f"{count} finding{'s' if count != 1 else ''}",
                 size=11, color=MUTED_ON_NAVY)
        top += 0.82
    # Trailing points share the same card chrome as the severity rows above.
    for point in list(spec.get("points") or [])[:2]:
        if top > 6.0:
            break
        _rect(slide, left, top, width, 0.7, NAVY_700, line_color=RULE_ON_NAVY)
        _rect(slide, left, top, 0.09, 0.7, GOLD)
        _textbox(slide, left + 0.25, top + 0.1, width - 0.5, 0.3,
                 safe_text(point.get("label")), size=12, color=GOLD, bold=True)
        _textbox(slide, left + 0.25, top + 0.38, width - 0.5, 0.28,
                 safe_text(point.get("text")), size=11, color=TEXT_ON_NAVY)
        top += 0.82


def _render_two_column(slide: Any, spec: dict[str, Any], body_top: float) -> None:
    body = spec.get("body")
    if body:
        _textbox(slide, MARGIN_IN, body_top, CONTENT_WIDTH_IN, 0.8, str(body),
                 size=16, color=TEXT_ON_NAVY)
    points = list(spec.get("points") or [])
    gap = 0.4
    column_width = (CONTENT_WIDTH_IN - gap) / 2
    split = (len(points) + 1) // 2
    columns = [points[:split], points[split:]]
    top = body_top + 1.0
    # Card height follows content: header room plus one row pitch per point.
    rows = max(1, min(4, max(len(column) for column in columns)))
    height = 0.5 + 0.82 * rows
    for index, column_points in enumerate(columns):
        left = MARGIN_IN + index * (column_width + gap)
        _rect(slide, left, top, column_width, height, NAVY_700,
              line_color=RULE_ON_NAVY)
        item_top = top + 0.25
        if not column_points:
            _textbox(slide, left + 0.25, item_top, column_width - 0.5, 0.4,
                     "No further points were supplied for this column.",
                     size=12, color=MUTED_ON_NAVY, italic=True)
            continue
        for point in column_points[:4]:
            _textbox(slide, left + 0.25, item_top, column_width - 0.5, 0.3,
                     safe_text(point.get("label")), size=13, color=GOLD, bold=True)
            _textbox(slide, left + 0.25, item_top + 0.32, column_width - 0.5, 0.5,
                     safe_text(point.get("text")), size=12, color=TEXT_ON_NAVY)
            item_top += 0.82


def _render_table(slide: Any, spec: dict[str, Any], data: dict[str, Any],
                  body_top: float) -> None:
    """Native table of the slide's points, or of the top findings when it has none."""
    body = spec.get("body")
    if body:
        _textbox(slide, MARGIN_IN, body_top, CONTENT_WIDTH_IN, 0.6, str(body),
                 size=16, color=TEXT_ON_NAVY)
    points = list(spec.get("points") or [])
    if points:
        headers = ["Signal", "What the evidence shows"]
        rows = [
            [safe_text(point.get("label")), safe_text(point.get("text"))]
            for point in points[:6]
        ]
        widths = (3.6, CONTENT_WIDTH_IN - 3.6)
    else:
        headers = ["Priority", "Finding", "Evidence"]
        rows = [
            [
                safe_text(finding.get("priority")),
                safe_text(finding.get("title")),
                ", ".join(str(item) for item in finding.get("evidence_ids") or [])
                or "Unavailable",
            ]
            for finding in list(data.get("findings") or [])[:6]
        ]
        widths = (1.5, CONTENT_WIDTH_IN - 4.5, 3.0)
    if not rows:
        _textbox(slide, MARGIN_IN, body_top + 1.0, CONTENT_WIDTH_IN, 0.6,
                 "No rows cleared evidence checks for this slide.",
                 size=15, color=MUTED_ON_NAVY, italic=True)
        return

    # Content-driven height, vertically centred in the zone under the body.
    row_height_in = 0.6
    zone_top = body_top + 0.8
    zone_bottom = 6.7
    table_height = min(zone_bottom - zone_top, row_height_in * (len(rows) + 1))
    table_top = zone_top + max(0.0, (zone_bottom - zone_top - table_height) / 2)
    shape = slide.shapes.add_table(
        len(rows) + 1, len(headers), Inches(MARGIN_IN), Inches(table_top),
        Inches(CONTENT_WIDTH_IN), Inches(table_height),
    )
    table = shape.table
    for index, width in enumerate(widths):
        table.columns[index].width = Emu(int(Inches(width)))
    for row in table.rows:
        row.height = Emu(int(Inches(row_height_in)))
    for column, header in enumerate(headers):
        cell = table.cell(0, column)
        cell.text = header
        cell.fill.solid()
        cell.fill.fore_color.rgb = _rgb(GOLD)
        paragraph = cell.text_frame.paragraphs[0]
        paragraph.runs[0].font.size = Pt(12)
        paragraph.runs[0].font.bold = True
        paragraph.runs[0].font.name = _FONTS.body
        paragraph.runs[0].font.color.rgb = _rgb(NAVY_900)
    for row_index, row_values in enumerate(rows, start=1):
        for column, value in enumerate(row_values):
            cell = table.cell(row_index, column)
            cell.text = value
            cell.fill.solid()
            cell.fill.fore_color.rgb = _rgb(NAVY_700 if row_index % 2 else NAVY_500)
            paragraph = cell.text_frame.paragraphs[0]
            paragraph.runs[0].font.size = Pt(11)
            paragraph.runs[0].font.name = _FONTS.body
            paragraph.runs[0].font.color.rgb = _rgb(TEXT_ON_NAVY)


def _render_statement(slide: Any, spec: dict[str, Any], body_top: float) -> None:
    body = spec.get("body")
    statement_top = body_top + 0.3
    if body:
        _textbox(slide, MARGIN_IN, statement_top, CONTENT_WIDTH_IN - 2.4,
                 3.95 - statement_top, str(body),
                 size=22, color=TEXT_ON_NAVY, font=_FONTS.display)
    _rect(slide, MARGIN_IN, 4.1, 2.2, 0.06, GOLD)
    points = list(spec.get("points") or [])[:3]
    top = 4.5
    if points:
        # A 3-across mini-card row (stat_rail style) in the 4.5–6.2in band:
        # the value is the hero (28pt gold) under an 11pt muted label.
        gap = 0.3
        card_width = (CONTENT_WIDTH_IN - 2 * gap) / 3
        card_height = 1.7
        for index, point in enumerate(points):
            left = MARGIN_IN + index * (card_width + gap)
            _rect(slide, left, top, card_width, card_height, NAVY_700,
                  line_color=RULE_ON_NAVY)
            _rect(slide, left, top, card_width, 0.07, GOLD)
            _textbox(slide, left + 0.22, top + 0.22, card_width - 0.44, 0.32,
                     safe_text(point.get("label")), size=11,
                     color=MUTED_ON_NAVY, bold=True)
            _textbox(slide, left + 0.22, top + 0.6, card_width - 0.44,
                     card_height - 0.75, safe_text(point.get("text")),
                     size=28, color=GOLD, bold=True, font=_FONTS.display)
        top += card_height
    callout = spec.get("callout")
    if callout:
        # Never collide with the card row: sit below whatever was drawn above.
        callout_top = max(6.35, top + 0.15)
        _rect(slide, MARGIN_IN, callout_top, CONTENT_WIDTH_IN, 0.55, NAVY_500,
              line_color=GOLD, line_width=1.0)
        _textbox(slide, MARGIN_IN + 0.2, callout_top + 0.09,
                 CONTENT_WIDTH_IN - 0.4, 0.4,
                 str(callout), size=13, color=GOLD_LIGHT, bold=True)


def _phase_summary(actions: list[dict[str, Any]]) -> list[dict[str, Any]]:
    phases: dict[str, dict[str, Any]] = {}
    order: list[str] = []
    for action in actions:
        phase = safe_text(action.get("phase"), "Plan")
        if phase not in phases:
            phases[phase] = {"name": phase, "weeks": [], "actions": []}
            order.append(phase)
        entry = phases[phase]
        week = action.get("week")
        if isinstance(week, int):
            entry["weeks"].append(week)
            week_end = action.get("week_end")
            if isinstance(week_end, int):
                entry["weeks"].append(week_end)
        entry["actions"].append(safe_text(action.get("action")))
    return [phases[name] for name in order[:5]]


def _render_timeline(slide: Any, spec: dict[str, Any], data: dict[str, Any],
                     body_top: float) -> None:
    """A true 16-week Gantt: one bar per phase across a week grid."""
    body = spec.get("body")
    if body:
        _textbox(slide, MARGIN_IN, body_top, CONTENT_WIDTH_IN, 0.5, str(body),
                 size=15, color=TEXT_ON_NAVY)
    phases = _phase_summary(list(data.get("actions") or []))
    if not phases:
        _textbox(
            slide, MARGIN_IN, body_top + 1.0, CONTENT_WIDTH_IN, 0.5,
            "The 16-week action plan is unavailable for this run.",
            size=16, color=MUTED_ON_NAVY,
        )
        return

    label_width = 2.9
    grid_left = MARGIN_IN + label_width
    grid_width = CONTENT_WIDTH_IN - label_width
    week_width = grid_width / 16
    header_top = body_top + 0.65
    for week in range(1, 17):
        left = grid_left + (week - 1) * week_width
        _rect(slide, left, header_top, week_width, 0.32,
              NAVY_700 if week % 2 else NAVY_900, line_color=RULE_ON_NAVY)
        _textbox(slide, left, header_top + 0.03, week_width, 0.26, f"W{week}",
                 size=9, color=MUTED_ON_NAVY, align=PP_ALIGN.CENTER)

    # Rows share the header-to-6.4in zone evenly so they fill it.
    row_top = header_top + 0.45
    row_pitch = (6.4 - row_top) / max(len(phases), 4)
    track_height = max(0.3, min(0.5, row_pitch - 0.28))
    for phase in phases:
        weeks = phase["weeks"]
        first = min(weeks) if weeks else None
        last = max(weeks) if weeks else None
        if first is None or last is None:
            week_label = "Weeks unscheduled"
            start, end = 1, 16
            bar_color, bar_text_color = NAVY_500, MUTED_ON_NAVY
        else:
            week_label = f"Week {first}" if first == last else f"Weeks {first}–{last}"
            start = max(1, min(16, first))
            end = max(start, min(16, last))
            bar_color, bar_text_color = GOLD, NAVY_900
        action_count = len(phase["actions"])
        actions_caption = f"{action_count} action{'s' if action_count != 1 else ''}"
        _textbox(slide, MARGIN_IN, row_top + 0.02, label_width - 0.15, 0.3,
                 phase["name"], size=13, color=TEXT_ON_NAVY, bold=True)
        if row_pitch >= 0.68:
            _textbox(slide, MARGIN_IN, row_top + 0.28, label_width - 0.15, 0.24,
                     week_label, size=10, color=MUTED_ON_NAVY)
            _textbox(slide, MARGIN_IN, row_top + 0.5, label_width - 0.15, 0.22,
                     actions_caption, size=9, color=MUTED_ON_NAVY)
        else:  # tight pitch: keep the label column to two lines
            _textbox(slide, MARGIN_IN, row_top + 0.26, label_width - 0.15, 0.22,
                     f"{week_label} · {actions_caption}", size=9, color=MUTED_ON_NAVY)
        _rect(slide, grid_left, row_top, grid_width, track_height, NAVY_700,
              line_color=RULE_ON_NAVY)
        bar_left = grid_left + (start - 1) * week_width
        bar_width = max(week_width * 0.6, (end - start + 1) * week_width)
        bar = _rect(slide, bar_left, row_top + 0.05, bar_width, track_height - 0.1,
                    bar_color, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
        frame = bar.text_frame
        frame.word_wrap = False
        for attr in ("margin_left", "margin_right", "margin_top", "margin_bottom"):
            setattr(frame, attr, Pt(2))
        # Never write action text inside a bar: only the week span, and only
        # when the bar is wide enough to carry it.
        if bar_width >= 1.2:
            run = frame.paragraphs[0].add_run()
            run.text = f"W{start}" if start == end else f"W{start}-{end}"
            run.font.size = Pt(9)
            run.font.bold = True
            run.font.name = _FONTS.body
            run.font.color.rgb = _rgb(bar_text_color)
        row_top += row_pitch

    _textbox(
        slide, MARGIN_IN, row_top + 0.05, CONTENT_WIDTH_IN, 0.32,
        "Bars show the scheduled week span of each phase in the canonical action plan.",
        size=10, color=MUTED_ON_NAVY, italic=True,
    )


def _render_comparison(slide: Any, spec: dict[str, Any], data: dict[str, Any],
                       body_top: float) -> None:
    """Clustered bar of client vs competitor medians from measured metrics."""
    body = spec.get("body")
    if body:
        _textbox(slide, MARGIN_IN, body_top, CONTENT_WIDTH_IN, 0.6, str(body),
                 size=15, color=TEXT_ON_NAVY)
    performance = data.get("performance_vs_competitors")
    performance = performance if isinstance(performance, dict) else {}
    metrics = [
        row for row in (performance.get("metrics") or [])
        if isinstance(row, dict)
        and isinstance(row.get("client"), int | float)
        and isinstance(row.get("competitor_median"), int | float)
    ]
    if metrics:
        chart_data = CategoryChartData()
        chart_data.categories = [safe_text(row.get("metric")) for row in metrics[:5]]
        chart_data.add_series(
            "Client", tuple(float(row["client"]) for row in metrics[:5])
        )
        chart_data.add_series(
            "Competitor median",
            tuple(float(row["competitor_median"]) for row in metrics[:5]),
        )
        graphic_frame = slide.shapes.add_chart(
            XL_CHART_TYPE.COLUMN_CLUSTERED,
            Inches(MARGIN_IN), Inches(body_top + 0.7),
            Inches(CONTENT_WIDTH_IN * 0.6), Inches(3.5), chart_data,
        )
        chart = graphic_frame.chart
        _dark_chart(chart)
        chart.has_title = False
        chart.has_legend = True
        chart.legend.position = XL_LEGEND_POSITION.BOTTOM
        chart.legend.include_in_layout = False
        chart.legend.font.size = Pt(11)
        chart.legend.font.name = _FONTS.body
        chart.legend.font.color.rgb = _rgb(TEXT_ON_NAVY)
        for index, series in enumerate(chart.plots[0].series):
            series.format.fill.solid()
            series.format.fill.fore_color.rgb = _rgb(
                GOLD if index == 0 else MUTED_ON_NAVY
            )
    else:
        _textbox(
            slide, MARGIN_IN, body_top + 0.9, CONTENT_WIDTH_IN * 0.6, 1.0,
            safe_text(
                performance.get("unavailable_reason"),
                "Competitor performance was not measured for this run, so no comparison "
                "is charted.",
            ),
            size=14, color=MUTED_ON_NAVY, italic=True,
        )

    left = MARGIN_IN + CONTENT_WIDTH_IN * 0.63
    width = CONTENT_WIDTH_IN - CONTENT_WIDTH_IN * 0.63
    top = body_top + 0.7
    summary = performance.get("summary")
    if summary:
        _textbox(slide, left, top, width, 1.0, str(summary), size=12,
                 color=TEXT_ON_NAVY)
        top += 1.15
    for row in (performance.get("metrics") or [])[:4]:
        if not isinstance(row, dict) or top > 6.0:
            continue
        _rect(slide, left, top, width, 0.72, NAVY_700, line_color=RULE_ON_NAVY)
        position = str(row.get("position") or "unknown").casefold()
        accent = {
            "ahead": POSITIVE, "behind": CRITICAL, "level": GOLD,
        }.get(position, RULE_ON_NAVY)
        _rect(slide, left, top, 0.08, 0.72, accent)
        _textbox(slide, left + 0.22, top + 0.08, width - 0.44, 0.3,
                 safe_text(row.get("metric")), size=12, color=TEXT_ON_NAVY, bold=True)
        _textbox(slide, left + 0.22, top + 0.38, width - 0.44, 0.3,
                 f"{safe_text(row.get('position'), 'unknown').title()} · "
                 f"{safe_text(row.get('note'), 'No note supplied')}",
                 size=10, color=MUTED_ON_NAVY)
        top += 0.82
    for point in list(spec.get("points") or [])[:2]:
        if top > 6.2:
            break
        _textbox(slide, left, top, width, 0.3, safe_text(point.get("label")),
                 size=12, color=GOLD, bold=True)
        _textbox(slide, left, top + 0.28, width, 0.4, safe_text(point.get("text")),
                 size=11, color=TEXT_ON_NAVY)
        top += 0.8


# --------------------------------------------------------------------------- theme


_MAJOR_FONT = re.compile(r'(<a:majorFont><a:latin typeface=")[^"]*')
_MINOR_FONT = re.compile(r'(<a:minorFont><a:latin typeface=")[^"]*')


def _patch_theme_fonts(presentation: Presentation) -> None:
    """Rewrite the theme's major/minor Latin typefaces to the brand pair."""
    try:
        theme_part = presentation.slide_master.part.part_related_by(RT.THEME)
        xml = theme_part.blob.decode("utf-8")
    except (KeyError, AttributeError, UnicodeDecodeError):  # pragma: no cover - defensive
        return
    xml = _MAJOR_FONT.sub(lambda m: m.group(1) + _FONTS.display, xml, count=1)
    xml = _MINOR_FONT.sub(lambda m: m.group(1) + _FONTS.body, xml, count=1)
    theme_part._blob = xml.encode("utf-8")


# --------------------------------------------------------------------------- entry


def render_deck(data: dict, output: Path) -> Path:
    """Render the executive deck from ``data['deck']`` slide specs."""
    output = Path(output)
    output.parent.mkdir(parents=True, exist_ok=True)

    presentation = Presentation()
    presentation.slide_width = Inches(SLIDE_WIDTH_IN)
    presentation.slide_height = Inches(SLIDE_HEIGHT_IN)
    _patch_theme_fonts(presentation)
    blank_layout = presentation.slide_layouts[6]

    specs = list(data.get("deck") or [])
    layouts = _resolve_layouts(specs)
    total = len(specs)
    for index, (spec, layout) in enumerate(zip(specs, layouts, strict=True), start=1):
        slide = presentation.slides.add_slide(blank_layout)
        if layout == "cover":
            _render_cover(slide, spec, data)
        else:
            _background(slide, NAVY_900)
            _eyebrow(slide, safe_text(spec.get("eyebrow"), ""))
            body_top = _title(slide, safe_text(spec.get("title"), "Untitled section"))
            if layout == "score":
                _render_score(slide, spec, data, body_top)
            elif layout == "stat_rail":
                _render_stat_rail(slide, spec, data, body_top)
            elif layout == "chart":
                _render_chart(slide, spec, data, body_top)
            elif layout == "two_column":
                _render_two_column(slide, spec, body_top)
            elif layout == "table":
                _render_table(slide, spec, data, body_top)
            elif layout == "timeline":
                _render_timeline(slide, spec, data, body_top)
            elif layout == "comparison":
                _render_comparison(slide, spec, data, body_top)
            else:
                _render_statement(slide, spec, body_top)
        # Every slide sits on navy now, so the mark always uses its dark form.
        _brand_mark(slide, dark=True)
        _footer(slide, data, index, total, dark=True)

    presentation.save(str(output))
    return output


def slide_layouts_for(data: dict) -> list[str]:
    """Expose the resolved layout sequence (used by tests and the PDF sibling)."""
    return _resolve_layouts(list(data.get("deck") or []))


# --------------------------------------------------------------------------- pdf sibling


def _pdf_fonts() -> tuple[str, str]:
    """Register the brand TTFs when they ship with the project; else fall back."""
    from reportlab.pdfbase import pdfmetrics
    from reportlab.pdfbase.ttfonts import TTFont

    from exporters.brand import font_paths

    display_path, body_path = font_paths(Path(__file__).resolve().parents[1])
    display, body = "Helvetica-Bold", "Helvetica"
    registered = set(pdfmetrics.getRegisteredFontNames())
    if display_path is not None:
        if _FONTS.display not in registered:
            pdfmetrics.registerFont(TTFont(_FONTS.display, str(display_path)))
        display = _FONTS.display
    if body_path is not None:
        name = _FONTS.body.replace(" ", "")
        if name not in registered:
            pdfmetrics.registerFont(TTFont(name, str(body_path)))
        body = name
    return display, body


def render_deck_pdf(data: dict, output: Path) -> Path:
    """Render the same slide sequence to PDF natively — one page per slide.

    ReportLab only: there is no LibreOffice on the host, so the PPTX is never
    converted. The page sequence, ordering and copy match ``render_deck``, and
    withheld scores stay withheld rather than being drawn as zero.
    """
    from reportlab.lib.colors import HexColor
    from reportlab.lib.pagesizes import A4, landscape
    from reportlab.lib.utils import simpleSplit
    from reportlab.pdfgen import canvas as pdf_canvas

    output = Path(output)
    output.parent.mkdir(parents=True, exist_ok=True)
    display_font, body_font = _pdf_fonts()
    width, height = landscape(A4)
    margin = 42.0

    specs = list(data.get("deck") or [])
    layouts = _resolve_layouts(specs)
    total = len(specs)
    client = safe_text(data.get("client", {}).get("name"), "Client")
    as_of = safe_text(data.get("run", {}).get("evidence_as_of"))

    document = pdf_canvas.Canvas(str(output), pagesize=(width, height))
    document.setTitle(f"{client} Enterprise SEO Executive Deck")
    document.setAuthor("Traffic Radius")
    document.setSubject("Evidence-led SEO approval deck")

    def text_block(text: str, x: float, y: float, font: str, size: int,
                   color: str, max_width: float, leading: float) -> float:
        document.setFont(font, size)
        document.setFillColor(HexColor(color))
        for line in simpleSplit(text, font, size, max_width):
            document.drawString(x, y, line)
            y -= leading
        return y

    def mark(x: float, y: float, dark: bool) -> None:
        for offset, bar_height, color in (
            (0.0, 8.0, WHITE if dark else BRAND_DEEP),
            (9.0, 13.0, BRAND_BLUE),
            (18.0, 18.0, BRAND_GREEN),
        ):
            document.setFillColor(HexColor(color))
            document.rect(x + offset, y, 6.0, bar_height, stroke=0, fill=1)

    for index, (spec, layout) in enumerate(zip(specs, layouts, strict=True), start=1):
        dark = layout == "cover"
        document.setFillColor(HexColor(NAVY_900))
        document.rect(0, 0, width, height, stroke=0, fill=1)
        if dark:
            document.setFillColor(HexColor(GOLD))
            document.rect(0, 0, 10, height, stroke=0, fill=1)

        content_width = width - 2 * margin
        cursor = height - margin - 24
        eyebrow = safe_text(spec.get("eyebrow"), "").upper()
        if eyebrow:
            cursor = text_block(
                eyebrow, margin, cursor, body_font, 10,
                GOLD, content_width, 14,
            ) - 8
        cursor = text_block(
            safe_text(spec.get("title"), "Untitled section"), margin, cursor,
            display_font, 26, TEXT_ON_NAVY, content_width, 32,
        ) - 10
        body = spec.get("body")
        if body:
            cursor = text_block(
                str(body), margin, cursor, body_font, 14,
                TEXT_ON_NAVY, content_width, 19,
            ) - 12

        if layout in {"score", "stat_rail"}:
            scored, withheld = _scored_categories(data)
            for category in scored:
                label = safe_text(category.get("category"))
                score = float(category["score"])
                document.setFillColor(HexColor(TEXT_ON_NAVY))
                document.setFont(body_font, 12)
                document.drawString(margin, cursor, label)
                bar_left = margin + 190
                bar_width = min(content_width - 260, 320.0)
                document.setFillColor(HexColor(RULE_ON_NAVY))
                document.rect(bar_left, cursor - 3, bar_width, 12, stroke=0, fill=1)
                document.setFillColor(HexColor(GOLD))
                document.rect(bar_left, cursor - 3,
                              bar_width * max(0.0, min(score, 100.0)) / 100.0, 12,
                              stroke=0, fill=1)
                document.setFillColor(HexColor(GOLD_LIGHT))
                document.drawString(bar_left + bar_width + 12, cursor, f"{score:.0f}")
                cursor -= 22
            for category in withheld:
                document.setFillColor(HexColor(TEXT_ON_NAVY))
                document.setFont(body_font, 12)
                document.drawString(margin, cursor, safe_text(category.get("category")))
                document.setFillColor(HexColor(MUTED_ON_NAVY))
                document.drawString(margin + 190, cursor, "Withheld")
                cursor -= 22
            reason = _withheld_reason(data, withheld)
            if reason:
                cursor = text_block(f"Withheld: {reason}", margin, cursor - 6,
                                    body_font, 10, MUTED_ON_NAVY, content_width, 13)
        elif layout == "timeline":
            for phase in _phase_summary(list(data.get("actions") or [])):
                weeks = phase["weeks"]
                span = (
                    f"Weeks {min(weeks)}–{max(weeks)}"
                    if weeks and min(weeks) != max(weeks)
                    else (f"Week {weeks[0]}" if weeks else "Weeks unscheduled")
                )
                cursor = text_block(f"{phase['name']} · {span}", margin, cursor,
                                    body_font, 12, TEXT_ON_NAVY, content_width, 16) - 4
        elif layout == "chart":
            for name, count in _severity_mix(data):
                cursor = text_block(
                    f"{name}: {count} finding{'s' if count != 1 else ''}",
                    margin, cursor, body_font, 12, TEXT_ON_NAVY, content_width, 16,
                )
        elif layout == "comparison":
            performance = data.get("performance_vs_competitors")
            performance = performance if isinstance(performance, dict) else {}
            rows = [row for row in (performance.get("metrics") or [])
                    if isinstance(row, dict)]
            if rows:
                for row in rows[:6]:
                    cursor = text_block(
                        f"{safe_text(row.get('metric'))} — client "
                        f"{safe_text(row.get('client'))}, competitor median "
                        f"{safe_text(row.get('competitor_median'))} "
                        f"({safe_text(row.get('position'), 'unknown')})",
                        margin, cursor, body_font, 11, TEXT_ON_NAVY, content_width, 15,
                    )
            else:
                cursor = text_block(
                    safe_text(
                        performance.get("unavailable_reason"),
                        "Competitor performance was not measured for this run.",
                    ),
                    margin, cursor, body_font, 12, MUTED_ON_NAVY, content_width, 16,
                )

        for point in list(spec.get("points") or [])[:4]:
            if cursor < margin + 60:
                break
            cursor = text_block(safe_text(point.get("label")), margin, cursor,
                                body_font, 12, GOLD, content_width, 15)
            cursor = text_block(safe_text(point.get("text")), margin, cursor,
                                body_font, 11, TEXT_ON_NAVY, content_width, 14) - 6

        # The navy page is dark on every slide, so the mark always uses its
        # dark form and the footer its muted-on-navy tone.
        mark(width - margin - 24, height - margin - 4, True)
        document.setFont(body_font, 8)
        document.setFillColor(HexColor(MUTED_ON_NAVY))
        document.drawString(
            margin, margin * 0.6,
            f"{client} · Enterprise SEO Audit · Evidence as of {as_of}",
        )
        document.drawRightString(
            width - margin, margin * 0.6, f"slide {index}/{total}"
        )
        document.showPage()

    if not specs:
        document.setFillColor(HexColor(NAVY_900))
        document.rect(0, 0, width, height, stroke=0, fill=1)
        text_block(
            "No deck slides were compiled for this run.", margin, height / 2,
            body_font, 14, MUTED_ON_NAVY, width - 2 * margin, 18,
        )
        document.showPage()

    document.save()
    return output
